"""
Content management tools for PowerPoint MCP Server.
Handles slides, text, images, and content manipulation.
"""
from typing import Dict, List, Optional, Any, Union
from mcp.server.fastmcp import FastMCP
import utils as ppt_utils
import tempfile
import base64
import os
import urllib.request
import urllib.parse

# Optional: requests for better HTTP handling
try:
    import requests  # type: ignore
except Exception:
    requests = None


def register_content_tools(app: FastMCP, presentations: Dict, get_current_presentation_id, validate_parameters, is_positive, is_non_negative, is_in_range, is_valid_rgb):
    """Register content management tools with the FastMCP app"""
    
    @app.tool()
    def add_slide(
        layout_index: int = 1,
        title: Optional[str] = None,
        background_type: Optional[str] = None,  # "solid", "gradient", "professional_gradient"
        background_colors: Optional[List[List[int]]] = None,  # For gradient: [[start_rgb], [end_rgb]]
        gradient_direction: str = "horizontal",
        color_scheme: str = "modern_blue",
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Add a new slide to the presentation with optional background styling."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        # Validate layout index
        if layout_index < 0 or layout_index >= len(pres.slide_layouts):
            return {
                "error": f"Invalid layout index: {layout_index}. Available layouts: 0-{len(pres.slide_layouts) - 1}"
            }
        
        try:
            # Add the slide
            slide, layout = ppt_utils.add_slide(pres, layout_index)
            slide_index = len(pres.slides) - 1
            
            # Set title if provided
            if title:
                ppt_utils.set_title(slide, title)
            
            # Apply background if specified
            if background_type == "gradient" and background_colors and len(background_colors) >= 2:
                ppt_utils.set_slide_gradient_background(
                    slide, background_colors[0], background_colors[1], gradient_direction
                )
            elif background_type == "professional_gradient":
                ppt_utils.create_professional_gradient_background(
                    slide, color_scheme, "subtle", gradient_direction
                )
            
            return {
                "message": f"Added slide {slide_index} with layout {layout_index}",
                "slide_index": slide_index,
                "layout_name": layout.name if hasattr(layout, 'name') else f"Layout {layout_index}"
            }
        except Exception as e:
            return {
                "error": f"Failed to add slide: {str(e)}"
            }

    @app.tool()
    def get_slide_info(slide_index: int, presentation_id: Optional[str] = None) -> Dict:
        """Get information about a specific slide."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        try:
            return ppt_utils.get_slide_info(slide, slide_index)
        except Exception as e:
            return {
                "error": f"Failed to get slide info: {str(e)}"
            }

    @app.tool()
    def extract_slide_text(slide_index: int, presentation_id: Optional[str] = None) -> Dict:
        """Extract all text content from a specific slide."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        try:
            result = ppt_utils.extract_slide_text_content(slide)
            result["slide_index"] = slide_index
            return result
        except Exception as e:
            return {
                "error": f"Failed to extract slide text: {str(e)}"
            }

    @app.tool()
    def extract_presentation_text(presentation_id: Optional[str] = None, include_slide_info: bool = True) -> Dict:
        """Extract all text content from all slides in the presentation."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        try:
            slides_text = []
            total_text_shapes = 0
            slides_with_tables = 0
            slides_with_titles = 0
            all_presentation_text = []
            
            for slide_index, slide in enumerate(pres.slides):
                slide_text_result = ppt_utils.extract_slide_text_content(slide)
                
                if slide_text_result["success"]:
                    slide_data = {
                        "slide_index": slide_index,
                        "text_content": slide_text_result["text_content"]
                    }
                    
                    if include_slide_info:
                        # Add basic slide info
                        slide_data["layout_name"] = slide.slide_layout.name
                        slide_data["total_text_shapes"] = slide_text_result["total_text_shapes"]
                        slide_data["has_title"] = slide_text_result["has_title"]
                        slide_data["has_tables"] = slide_text_result["has_tables"]
                    
                    slides_text.append(slide_data)
                    
                    # Accumulate statistics
                    total_text_shapes += slide_text_result["total_text_shapes"]
                    if slide_text_result["has_tables"]:
                        slides_with_tables += 1
                    if slide_text_result["has_title"]:
                        slides_with_titles += 1
                    
                    # Collect all text for combined output
                    if slide_text_result["text_content"]["all_text_combined"]:
                        all_presentation_text.append(f"=== SLIDE {slide_index + 1} ===")
                        all_presentation_text.append(slide_text_result["text_content"]["all_text_combined"])
                        all_presentation_text.append("")  # Empty line separator
                else:
                    slides_text.append({
                        "slide_index": slide_index,
                        "error": slide_text_result.get("error", "Unknown error"),
                        "text_content": None
                    })
            
            return {
                "success": True,
                "presentation_id": pres_id,
                "total_slides": len(pres.slides),
                "slides_with_text": len([s for s in slides_text if s.get("text_content") is not None]),
                "total_text_shapes": total_text_shapes,
                "slides_with_titles": slides_with_titles,
                "slides_with_tables": slides_with_tables,
                "slides_text": slides_text,
                "all_presentation_text_combined": "\n".join(all_presentation_text)
            }
            
        except Exception as e:
            return {
                "error": f"Failed to extract presentation text: {str(e)}"
            }

    @app.tool()
    def populate_placeholder(
        slide_index: int,
        placeholder_idx: int,
        text: str,
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Populate a placeholder with text."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        try:
            ppt_utils.populate_placeholder(slide, placeholder_idx, text)
            return {
                "message": f"Populated placeholder {placeholder_idx} on slide {slide_index}"
            }
        except Exception as e:
            return {
                "error": f"Failed to populate placeholder: {str(e)}"
            }

    @app.tool()
    def add_bullet_points(
        slide_index: int,
        placeholder_idx: int,
        bullet_points: List[str],
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Add bullet points to a placeholder."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        try:
            placeholder = slide.placeholders[placeholder_idx]
            ppt_utils.add_bullet_points(placeholder, bullet_points)
            return {
                "message": f"Added {len(bullet_points)} bullet points to placeholder {placeholder_idx} on slide {slide_index}"
            }
        except Exception as e:
            return {
                "error": f"Failed to add bullet points: {str(e)}"
            }

    @app.tool()
    def manage_text(
        slide_index: int,
        operation: str,  # "add", "format", "validate", "format_runs"
        left: float = 1.0,
        top: float = 1.0,
        width: float = 4.0,
        height: float = 2.0,
        text: str = "",
        shape_index: Optional[int] = None,  # For format/validate operations
        text_runs: Optional[List[Dict]] = None,  # For format_runs operation
        # Formatting options
        font_size: Optional[int] = None,
        font_name: Optional[str] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        underline: Optional[bool] = None,
        color: Optional[List[int]] = None,
        bg_color: Optional[List[int]] = None,
        alignment: Optional[str] = None,
        vertical_alignment: Optional[str] = None,
        # Advanced options
        auto_fit: bool = True,
        validation_only: bool = False,
        min_font_size: int = 8,
        max_font_size: int = 72,
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Unified text management tool for adding, formatting, validating text, and formatting multiple text runs."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        # Validate parameters
        validations = {}
        if font_size is not None:
            validations["font_size"] = (font_size, [(is_positive, "must be a positive integer")])
        if color is not None:
            validations["color"] = (color, [(is_valid_rgb, "must be a valid RGB list [R, G, B] with values 0-255")])
        if bg_color is not None:
            validations["bg_color"] = (bg_color, [(is_valid_rgb, "must be a valid RGB list [R, G, B] with values 0-255")])
        
        if validations:
            valid, error = validate_parameters(validations)
            if not valid:
                return {"error": error}
        
        try:
            if operation == "add":
                # Auto-detect URL even if source_type is not explicitly "url"
                if isinstance(image_source, str) and (image_source.startswith("http://") or image_source.startswith("https://")):
                    source_type = "url"
                # Add new textbox
                shape = ppt_utils.add_textbox(
                    slide, left, top, width, height, text,
                    font_size=font_size,
                    font_name=font_name,
                    bold=bold,
                    italic=italic,
                    underline=underline,
                    color=tuple(color) if color else None,
                    bg_color=tuple(bg_color) if bg_color else None,
                    alignment=alignment,
                    vertical_alignment=vertical_alignment,
                    auto_fit=auto_fit
                )
                return {
                    "message": f"Added text box to slide {slide_index}",
                    "shape_index": len(slide.shapes) - 1,
                    "text": text
                }
            
            elif operation == "format":
                # Format existing text shape
                if shape_index is None or shape_index < 0 or shape_index >= len(slide.shapes):
                    return {
                        "error": f"Invalid shape index for formatting: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
                    }
                
                shape = slide.shapes[shape_index]
                ppt_utils.format_text_advanced(
                    shape,
                    font_size=font_size,
                    font_name=font_name,
                    bold=bold,
                    italic=italic,
                    underline=underline,
                    color=tuple(color) if color else None,
                    bg_color=tuple(bg_color) if bg_color else None,
                    alignment=alignment,
                    vertical_alignment=vertical_alignment
                )
                return {
                    "message": f"Formatted text shape {shape_index} on slide {slide_index}"
                }
            
            elif operation == "validate":
                # Validate text fit
                if shape_index is None or shape_index < 0 or shape_index >= len(slide.shapes):
                    return {
                        "error": f"Invalid shape index for validation: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
                    }
                
                validation_result = ppt_utils.validate_text_fit(
                    slide.shapes[shape_index],
                    text_content=text or None,
                    font_size=font_size or 12
                )
                
                if not validation_only and validation_result.get("needs_optimization"):
                    # Apply automatic fixes
                    fix_result = ppt_utils.validate_and_fix_slide(
                        slide,
                        auto_fix=True,
                        min_font_size=min_font_size,
                        max_font_size=max_font_size
                    )
                    validation_result.update(fix_result)
                
                return validation_result
            
            elif operation == "format_runs":
                # Format multiple text runs with different formatting
                if shape_index is None or shape_index < 0 or shape_index >= len(slide.shapes):
                    return {
                        "error": f"Invalid shape index for format_runs: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
                    }
                
                if not text_runs:
                    return {"error": "text_runs parameter is required for format_runs operation"}
                
                shape = slide.shapes[shape_index]
                
                # Check if shape has text
                if not hasattr(shape, 'text_frame') or not shape.text_frame:
                    return {"error": "Shape does not contain text"}
                
                # Clear existing text and rebuild with formatted runs
                text_frame = shape.text_frame
                text_frame.clear()
                
                formatted_runs = []
                
                for run_data in text_runs:
                    if 'text' not in run_data:
                        continue
                        
                    # Add paragraph if needed
                    if not text_frame.paragraphs:
                        paragraph = text_frame.paragraphs[0]
                    else:
                        paragraph = text_frame.add_paragraph()
                    
                    # Add run with text
                    run = paragraph.add_run()
                    run.text = run_data['text']
                    
                    # Apply formatting using pptx imports
                    from pptx.util import Pt
                    from pptx.dml.color import RGBColor
                    
                    if 'bold' in run_data:
                        run.font.bold = run_data['bold']
                    if 'italic' in run_data:
                        run.font.italic = run_data['italic']
                    if 'underline' in run_data:
                        run.font.underline = run_data['underline']
                    if 'font_size' in run_data:
                        run.font.size = Pt(run_data['font_size'])
                    if 'font_name' in run_data:
                        run.font.name = run_data['font_name']
                    if 'color' in run_data and is_valid_rgb(run_data['color']):
                        run.font.color.rgb = RGBColor(*run_data['color'])
                    if 'hyperlink' in run_data:
                        run.hyperlink.address = run_data['hyperlink']
                    
                    formatted_runs.append({
                        "text": run_data['text'],
                        "formatting_applied": {k: v for k, v in run_data.items() if k != 'text'}
                    })
                
                return {
                    "message": f"Applied formatting to {len(formatted_runs)} text runs on shape {shape_index}",
                    "slide_index": slide_index,
                    "shape_index": shape_index,
                    "formatted_runs": formatted_runs
                }
            
            else:
                return {
                    "error": f"Invalid operation: {operation}. Must be 'add', 'format', 'validate', or 'format_runs'"
                }
        
        except Exception as e:
            return {
                "error": f"Failed to {operation} text: {str(e)}"
            }

    @app.tool()
    def manage_image(
        slide_index: int,
        operation: str,  # "add", "enhance"
        image_source: str,  # file path or base64 string
        source_type: str = "file",  # "file" or "base64"
        left: float = 1.0,
        top: float = 1.0,
        width: Optional[float] = None,
        height: Optional[float] = None,
        # Enhancement options
        enhancement_style: Optional[str] = None,  # "presentation", "custom"
        brightness: float = 1.0,
        contrast: float = 1.0,
        saturation: float = 1.0,
        sharpness: float = 1.0,
        blur_radius: float = 0,
        filter_type: Optional[str] = None,
        output_path: Optional[str] = None,
        presentation_id: Optional[str] = None
    ) -> Dict:
        """
        统一的图片处理工具（添加/增强）。

        功能
        - operation="add"：将图片插入到指定幻灯片位置，支持本地文件或 Base64 图片源或图片地址。
        - operation="enhance"：对已有图片文件进行画质增强与风格化处理，输出增强后的图片路径。

        参数
        - slide_index: int — 目标幻灯片索引（从 0 开始）。
        - operation: str — "add" 或 "enhance"。
        - image_source: str —
          * 当 source_type="file"：本地图片文件路径。
          * 当 source_type="base64"：图片的 Base64 字符串。
          * 当 source_type="url"：图片的 http/https 地址。
        - source_type: str — "file"、"base64" 或 "url"。
          * add 支持 "file"、"base64"、"url"（仅允许 http/https）。
          * enhance 仅支持 "file"（不接受 base64 或 url）。
        - left, top: float — 插入位置（英寸）。
        - width, height: Optional[float] — 插入尺寸（英寸）。可只提供一项以按比例缩放；都不提供则按图片原始尺寸。
        - enhancement_style: Optional[str] — "presentation" 或 "custom"。当 operation="add" 且需要自动增强时可用；"presentation" 走预设的专业增强流程。
        - brightness, contrast, saturation, sharpness: float — 亮度/对比度/饱和度/锐度（默认 1.0，>1 增强，<1 减弱）。
        - blur_radius: float — 模糊半径（默认 0）。
        - filter_type: Optional[str] — 过滤器类型（如 "DETAIL"、"SMOOTH" 等，取决于 Pillow 支持）。
        - output_path: Optional[str] — 增强后图片的输出路径（不传则生成临时文件）。
        - presentation_id: Optional[str] — 指定演示文稿 ID；不传则使用当前打开的演示文稿。
          注：在某些部署环境（如你当前环境）中，必须显式传入 presentation_id 才会对目标文档生效。

        返回
        - operation="add"：
          * message: str
          * shape_index: int — 新增形状索引
          * image_path: str（当 source_type="file" 时返回）
        - operation="enhance"：
          * message: str
          * enhanced_path: str — 增强后图片文件路径
        - 失败时返回 {"error": str}

        注意事项
        - 现在支持通过 URL 插入图片（仅 operation="add"）：source_type="url"，仅允许 http/https 协议。
          会在内部下载到临时文件后插入并自动清理。若返回的 Content-Type 非 image/* 将返回错误。
          enhance 仍然仅支持本地文件路径。
        - 在某些部署环境（如你当前环境）中，需显式提供 presentation_id 参数，否则可能插入到非预期文档或不生效。
        - operation="enhance" 不接受 Base64，必须提供可访问的本地文件路径。
        - 插入 Base64 图片时，内部会写入临时文件后再插入，操作完成后临时文件会被清理。
        - slide_index 必须在当前演示文稿的有效范围内，否则将返回错误。

        示例
        - 通过 URL 插入图片（仅 add）：
            manage_image(slide_index=0, operation="add",
                         image_source="https://example.com/logo.png", source_type="url",
                         left=1.0, top=1.0, width=3.0, presentation_id="YOUR_PRESENTATION_ID")
        - 插入本地图片：
            manage_image(slide_index=0, operation="add",
                         image_source="D:/images/logo.png", source_type="file",
                         left=1.0, top=1.0, width=3.0, presentation_id="YOUR_PRESENTATION_ID")

        - 插入 Base64 图片：
            manage_image(slide_index=1, operation="add",
                         image_source="<BASE64字符串>", source_type="base64",
                         left=2.0, top=1.5, width=4.0, height=2.5, presentation_id="YOUR_PRESENTATION_ID")

        - 插入并应用专业增强（演示风格）：
            manage_image(slide_index=2, operation="add",
                         image_source="assets/photo.jpg", source_type="file",
                         enhancement_style="presentation", left=1.0, top=2.0, presentation_id="YOUR_PRESENTATION_ID")

        - 增强已有图片文件（自定义参数）：
            manage_image(slide_index=0, operation="enhance",
                         image_source="assets/photo.jpg", source_type="file",
                         brightness=1.2, contrast=1.1, saturation=1.3,
                         sharpness=1.1, blur_radius=0, filter_type=None,
                         output_path="assets/photo_enhanced.jpg", presentation_id="YOUR_PRESENTATION_ID")
        """
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        try:
            if operation == "add":
                if source_type == "base64":
                    # Handle base64 image
                    try:
                        image_data = base64.b64decode(image_source)
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_file:
                            temp_file.write(image_data)
                            temp_path = temp_file.name
                        
                        # Add image from temporary file
                        shape = ppt_utils.add_image(slide, temp_path, left, top, width, height)
                        
                        # Clean up temporary file
                        os.unlink(temp_path)
                        
                        return {
                            "message": f"Added image from base64 to slide {slide_index}",
                            "shape_index": len(slide.shapes) - 1
                        }
                    except Exception as e:
                        return {
                            "error": f"Failed to process base64 image: {str(e)}"
                        }
                elif source_type == "url":
                    # Handle image URL (http/https)
                    try:
                        # Normalize and percent-encode URL path/query to support spaces and non-ASCII characters
                        parsed = urllib.parse.urlsplit(image_source)
                        if parsed.scheme not in ("http", "https"):
                            return {"error": f"Unsupported URL scheme: {parsed.scheme}. Only http/https allowed."}
                        encoded_path = urllib.parse.quote(parsed.path or "", safe="/%")
                        # Re-encode query preserving keys and multiple values
                        qsl = urllib.parse.parse_qsl(parsed.query or "", keep_blank_values=True)
                        encoded_query = urllib.parse.urlencode(qsl, doseq=True)
                        encoded_url = urllib.parse.urlunsplit((parsed.scheme, parsed.netloc, encoded_path, encoded_query, parsed.fragment))
                        
                        # Download helper using requests if available, else urllib
                        content_type = None
                        temp_path = None
                        image_exts = (".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tif", ".tiff")
                        
                        if requests is not None:
                            with requests.get(encoded_url, stream=True) as resp:
                                if resp.status_code != 200:
                                    return {"error": f"Failed to download image. HTTP {resp.status_code}"}
                                content_type = resp.headers.get("Content-Type", "") or ""
                                
                                # Determine suffix and allow fallback by URL extension if Content-Type missing or not image/*
                                suffix = ".png"
                                is_image = content_type.startswith("image/")
                                try:
                                    main_type = (content_type.split(";")[0].strip() if content_type else "")
                                    if "/" in main_type:
                                        ext = main_type.split("/")[1].lower()
                                        if ext in ("jpeg", "pjpeg"):
                                            suffix = ".jpg"
                                        elif ext in ("png", "gif", "bmp", "webp", "tiff"):
                                            suffix = f".{ext}"
                                except Exception:
                                    pass
                                if suffix == ".png":
                                    path_ext = os.path.splitext(parsed.path or "")[1].lower()
                                    if path_ext in image_exts:
                                        suffix = path_ext
                                
                                if not is_image and suffix not in image_exts:
                                    return {"error": f"URL content is not an image (Content-Type: {content_type or 'unknown'})"}
                                
                                with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
                                    temp_path = temp_file.name
                                    for chunk in resp.iter_content(chunk_size=8192):
                                        if not chunk:
                                            continue
                                        temp_file.write(chunk)
                        else:
                            req = urllib.request.Request(encoded_url, headers={"User-Agent": "Mozilla/5.0"})
                            with urllib.request.urlopen(req) as resp:
                                content_type = resp.headers.get("Content-Type", "") or ""
                                
                                suffix = ".png"
                                is_image = content_type.startswith("image/")
                                try:
                                    main_type = (content_type.split(";")[0].strip() if content_type else "")
                                    if "/" in main_type:
                                        ext = main_type.split("/")[1].lower()
                                        if ext in ("jpeg", "pjpeg"):
                                            suffix = ".jpg"
                                        elif ext in ("png", "gif", "bmp", "webp", "tiff"):
                                            suffix = f".{ext}"
                                except Exception:
                                    pass
                                if suffix == ".png":
                                    path_ext = os.path.splitext(parsed.path or "")[1].lower()
                                    if path_ext in image_exts:
                                        suffix = path_ext
                                
                                if not is_image and suffix not in image_exts:
                                    return {"error": f"URL content is not an image (Content-Type: {content_type or 'unknown'})"}
                                
                                with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
                                    temp_path = temp_file.name
                                    while True:
                                        chunk = resp.read(8192)
                                        if not chunk:
                                            break
                                        temp_file.write(chunk)
                        
                        # Add image from temporary file
                        shape = ppt_utils.add_image(slide, temp_path, left, top, width, height)
                        
                        # Clean up temporary file
                        if temp_path and os.path.exists(temp_path):
                            os.unlink(temp_path)
                        
                        return {
                            "message": f"Added image from URL to slide {slide_index}",
                            "shape_index": len(slide.shapes) - 1
                        }
                    except Exception as e:
                        # Best-effort cleanup if temp_path was created
                        try:
                            if temp_path and os.path.exists(temp_path):
                                os.unlink(temp_path)
                        except Exception:
                            pass
                        return {"error": f"Failed to process image URL: {str(e)}"}
                else:
                    # Handle file path
                    if not os.path.exists(image_source):
                        return {
                            "error": f"Image file not found: {image_source}"
                        }
                    
                    shape = ppt_utils.add_image(slide, image_source, left, top, width, height)
                    return {
                        "message": f"Added image to slide {slide_index}",
                        "shape_index": len(slide.shapes) - 1,
                        "image_path": image_source
                    }
            
            elif operation == "enhance":
                # Enhance existing image file
                if source_type == "base64":
                    return {
                        "error": "Enhancement operation requires file path, not base64 data"
                    }
                
                if not os.path.exists(image_source):
                    return {
                        "error": f"Image file not found: {image_source}"
                    }
                
                if enhancement_style == "presentation":
                    # Apply professional enhancement
                    enhanced_path = ppt_utils.apply_professional_image_enhancement(
                        image_source, style="presentation", output_path=output_path
                    )
                else:
                    # Apply custom enhancement
                    enhanced_path = ppt_utils.enhance_image_with_pillow(
                        image_source,
                        brightness=brightness,
                        contrast=contrast,
                        saturation=saturation,
                        sharpness=sharpness,
                        blur_radius=blur_radius,
                        filter_type=filter_type,
                        output_path=output_path
                    )
                
                return {
                    "message": f"Enhanced image: {image_source}",
                    "enhanced_path": enhanced_path
                }
            
            else:
                return {
                    "error": f"Invalid operation: {operation}. Must be 'add' or 'enhance'"
                }
        
        except Exception as e:
            return {
                "error": f"Failed to {operation} image: {str(e)}"
            }