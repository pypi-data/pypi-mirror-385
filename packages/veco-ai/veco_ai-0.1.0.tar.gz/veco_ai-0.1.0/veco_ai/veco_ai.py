# veco_ai.py
# -----------------------------------------------------------------------------
# veco-ai - Stand-alone vectorizer and RAG retrieval toolkit
#
# - JSON fallback storage (no additional service required)
# - Optional external storage (SQLite/Mongo) via the separate "storages.py" module
# - Fixes/Features:
#     * Correct FAISS IDs (no more zero IDs)
#     * Overlapping chunking (~500-700 tokens, character based, sentence aware)
#     * Summaries are stored separately (never used as embedding input)
#     * RAG query helper with Ollama: query(database, question, llm_model, ...)
#     * Consistent relative path handling (no user specific absolute paths)
#     * OPTIONAL: Speaker diarization via the external veco_diarization.py module
#     * OPTIONAL: CNN based image classification (torchvision) and optional captions via veco_pic_describe
#     * AUTO heuristics choose the most sensible pipeline per file type when nothing is explicitly configured
# -----------------------------------------------------------------------------

from __future__ import annotations
import os
import sys
import time
import json
import hashlib
import logging
import threading
import importlib
import tempfile
from pathlib import Path
from typing import List, Dict, Any, Optional

import numpy as np

try:
    import torch  # type: ignore
except ImportError as exc:  # pragma: no cover - import guard
    torch = None  # type: ignore
    _TORCH_IMPORT_ERROR = exc
else:
    _TORCH_IMPORT_ERROR = None

try:
    import whisper  # type: ignore
except ImportError as exc:  # pragma: no cover - import guard
    whisper = None  # type: ignore
    _WHISPER_IMPORT_ERROR = exc
else:
    _WHISPER_IMPORT_ERROR = None

try:
    from sentence_transformers import SentenceTransformer
except ImportError as exc:  # pragma: no cover - import guard
    SentenceTransformer = None  # type: ignore
    _SENTENCE_TRANSFORMERS_IMPORT_ERROR = exc
else:
    _SENTENCE_TRANSFORMERS_IMPORT_ERROR = None

try:
    from faiss import IndexFlatL2, IndexIDMap
except ImportError as exc:  # pragma: no cover - import guard
    IndexFlatL2 = IndexIDMap = None  # type: ignore
    _FAISS_IMPORT_ERROR = exc
else:
    _FAISS_IMPORT_ERROR = None

import pdfplumber
import pytesseract
from PIL import Image
from docx import Document
from pptx import Presentation

try:
    from moviepy.editor import VideoFileClip
except Exception as exc:  # pragma: no cover - import guard
    VideoFileClip = None  # type: ignore
    _MOVIEPY_IMPORT_ERROR = exc
else:
    _MOVIEPY_IMPORT_ERROR = None


def _require_dependency(module, name: str, import_error: Exception | None):
    """Raise a clear error if a mandatory dependency is missing."""
    if module is None:
        message = (
            f"The library '{name}' is required for this feature. "
            "Please install the dependency as documented in the README/requirements."
        )
        if import_error is not None:
            raise RuntimeError(message) from import_error
        raise RuntimeError(message)

# Optional: Ollama integration (used only when available)
try:
    import ollama
except Exception:
    ollama = None

# Optional: Vision - torchvision classification (ResNet50)
try:
    import torchvision
    from torchvision import transforms
    _VISION_OK = True
except Exception:
    _VISION_OK = False
    torchvision = None
    transforms = None

# ---------------------------------- Logging ----------------------------------
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger("veco")


# -------------------------- Simple Embeddings -------------------------
class FallbackSentenceEmbedder:
    """Deterministic hashing-based embedder used when SBERT is unavailable."""

    def __init__(self, dim: int = 384):
        self.dim = dim

    def get_sentence_embedding_dimension(self) -> int:
        return self.dim

    def encode(self, texts: List[str], convert_to_numpy: bool = True, normalize_embeddings: bool = False) -> np.ndarray:
        vectors = np.zeros((len(texts), self.dim), dtype=np.float32)
        if not texts:
            return vectors
        for idx, text in enumerate(texts):
            digest = hashlib.sha256((text or "").encode("utf-8", errors="ignore")).digest()
            expanded = np.frombuffer(digest, dtype=np.uint8).astype(np.float32)
            repeats = int(np.ceil(self.dim / expanded.size))
            tiled = np.tile(expanded, repeats)[: self.dim]
            vec = tiled / 255.0
            if normalize_embeddings:
                norm = np.linalg.norm(vec)
                if norm > 0:
                    vec = vec / norm
            vectors[idx] = vec
        return vectors


# ---------------------------------- Spinner ----------------------------------
class Spinner:
    """Small CLI spinner for longer operations (purely cosmetic)."""
    spinner_cycle = ["|", "/", "-", "\\"]

    def __init__(self, message: str = "Processing"):
        self.stop_running = False
        self.thread = threading.Thread(target=self._run, daemon=True)
        self.message = message

    def _run(self):
        i = 0
        while not self.stop_running:
            sys.stdout.write(f"\r{self.message} {self.spinner_cycle[i % len(self.spinner_cycle)]}")
            sys.stdout.flush()
            time.sleep(0.08)
            i += 1
        sys.stdout.write("\r" + " " * (len(self.message) + 2) + "\r")

    def start(self):
        try:
            self.thread.start()
        except RuntimeError:
            pass

    def stop(self):
        self.stop_running = True
        try:
            self.thread.join(timeout=0.2)
        except RuntimeError:
            pass


# ------------------------------ Helpers / Utilities ------------------------------
def _try_import_storages():
    """Lazy-Import eines optionalen Moduls 'storages.py' (SQLite/Mongo-Backends)."""
    try:
        return importlib.import_module("storages")
    except Exception:
        return None


def _try_import_diarization():
    """Lazy-Import eines optionalen Moduls 'veco_diarization.py' (Speaker Diarization)."""
    try:
        return importlib.import_module("veco_diarization")
    except Exception:
        return None


def _try_import_pic_describe():
    """Lazy-Import eines optionalen Moduls 'veco_pic_describe' (externe Bildbeschreibung)."""
    try:
        return importlib.import_module("veco_pic_describe")
    except Exception:
        return None


def _relpath(p: str) -> str:
    """Return path p relative to the current working directory whenever possible."""
    try:
        if os.path.isabs(p):
            return os.path.relpath(p, start=os.getcwd())
        return p
    except Exception:
        return p  # Return unchanged if anything goes wrong


def chunk_text(text: str, chunk_chars: int = 1800, overlap_chars: int = 200) -> List[str]:
    """
    Character based chunking with overlap (no extra dependencies).
    - Versucht, bevorzugt an Satzenden ('.') zu schneiden.
    - Rule of thumb: 1800 characters roughly equals 500-700 tokens (depending on language/text).
    """
    text = (text or "").strip()
    if not text:
        return []
    chunks: List[str] = []
    n = len(text)
    i = 0
    while i < n:
        end = min(i + chunk_chars, n)
        cut = text.rfind(".", i, end)
        if cut == -1 or cut < i + int(0.6 * chunk_chars):
            cut = end
        chunk = text[i:cut].strip()
        if chunk:
            chunks.append(chunk)
        if cut >= n:
            break
        i = max(0, cut - overlap_chars)
    return chunks


# --------------------------- Main Class: Vectorize ---------------------------
class Vectorize:
    """
    - Loads/initialises embedding and ASR models
    - Extracts text from files (txt/pdf/docx/pptx/image/audio/video)
    - Chunks, embeds, and indexes content (FAISS)
    - Stores data in JSON (fallback) or optional backends (SQLite/MongoDB)
    - OPTIONAL: Speaker diarization (external pipeline), image classification (torchvision),
                image captioning (external veco_pic_describe module)
    - Provides RAG retrieval plus Ollama-powered querying
    """

    def __init__(
        self,
        default_model: str = "gemma3:12b",
        preload_json_path: Optional[str] = "vector_db.json",
        storage: Optional[object] = None,
        storage_kind: Optional[str] = None,
        storage_kwargs: Optional[dict] = None,
        write_through: bool = True,
        enable_audio: bool = True,
        audio_model_size: str = "base",
        whisper_download_root: Optional[str] = None,
        fallback_embedding_dim: int = 384,
        force_fallback_embedder: bool = False,
    ):
        _require_dependency(torch, "torch", _TORCH_IMPORT_ERROR)
        _require_dependency(whisper, "openai-whisper", _WHISPER_IMPORT_ERROR)
        _require_dependency(SentenceTransformer, "sentence-transformers", _SENTENCE_TRANSFORMERS_IMPORT_ERROR)
        _require_dependency(IndexFlatL2, "faiss-cpu", _FAISS_IMPORT_ERROR)
        _require_dependency(IndexIDMap, "faiss-cpu", _FAISS_IMPORT_ERROR)

        # Base configuration
        self.default_model = default_model
        self.preload_json_path = _relpath(preload_json_path or "vector_db.json")
        self.write_through = write_through

        # Internal storage (for the JSON fallback and quick save/load)
        self.outputdb: List[Dict[str, Any]] = []
        self.id_lookup: Dict[int, Dict[str, Any]] = {}
        self._next_vector_id = 0
        self._next_doc_id = 0
        self._active_db: Optional[str] = None

        # Optional external storage (SQLite/Mongo)
        self._ext_storage = None
        if storage is not None:
            self._ext_storage = storage
        elif storage_kind is not None:
            _stor = _try_import_storages()
            if _stor is not None:
                if storage_kind.lower() == "sqlite":
                    self._ext_storage = _stor.SqliteStorage(**(storage_kwargs or {}))
                elif storage_kind.lower() == "mongo":
                    self._ext_storage = _stor.MongoStorage(**(storage_kwargs or {}))
                else:
                    raise ValueError(f"Unbekanntes storage_kind: {storage_kind}")
            else:
                logger.warning("storages.py not found - staying with the JSON fallback.")

        # Load models (Whisper + SBERT, optionally check Ollama)
        self._audio_requested = enable_audio
        self._audio_available = False
        self.whisper_model = None
        self.audio_model_size = audio_model_size
        self._fallback_embedder: Optional[FallbackSentenceEmbedder] = None

        spinner = Spinner("Initializing models")
        spinner.start()
        try:
            if self._audio_requested:
                device = "cuda" if torch.cuda.is_available() else "cpu"
                try:
                    self.whisper_model = whisper.load_model(audio_model_size, device=device)
                except Exception as exc:
                    logger.warning("Whisper initialisation failed (%s). Audio disabled.", exc)
                    self.whisper_model = None
                    self._audio_available = False
                else:
                    self._audio_available = True
            else:
                self.whisper_model = None
                self._audio_available = False
            self.embedder = SentenceTransformer("all-MiniLM-L6-v2")
            self._embedding_dim = int(self.embedder.get_sentence_embedding_dimension())
            self.faiss_index = IndexIDMap(IndexFlatL2(self._embedding_dim))
            if ollama is not None:
                self.check_ollama_models()
        finally:
            spinner.stop()

        # Vision (optional) - lazy init handles
        self._vision_cls = None  # torchvision resnet50
        self._vision_tf = None   # associated transforms

        # Bootstrap: load existing data (prefer storage, otherwise JSON)
        if self._ext_storage is not None:
            self._bootstrap_from_storage()
        else:
            self.load_database()

    # ---------------------- Infrastructure / I/O ------------------------
    def check_ollama_models(self):
        """Only checks whether Ollama is reachable - never raises fatal errors."""
        try:
            _ = ollama.list()
        except Exception:
            logger.info("Ollama not available - LLM compression/RAG answer disabled.")

    @property
    def audio_available(self) -> bool:
        """Indicates if Whisper-based transcription can be used."""
        return bool(self._audio_available)

    def detect_input_type(self, path: str) -> str:
        """Dateityp heuristisch per Endung bestimmen."""
        p = str(path).lower()
        if p.endswith(".txt"):
            return "text"
        if p.endswith(".pdf"):
            return "pdf"
        if p.endswith(".docx"):
            return "word"
        if p.endswith(".pptx"):
            return "pptx"
        if p.endswith((".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff")):
            return "image"
        if p.endswith((".wav", ".mp3", ".m4a", ".flac", ".ogg")):
            return "audio"
        if p.endswith((".mp4", ".mov", ".mkv", ".avi")):
            return "video"
        return "text"

    # --------------------------- Extraction ----------------------------
    def extract_text(self, inputfile: str, input_type: str) -> str:
        """Extrahiert Text aus Text/PDF/Word/PPTX-Dateien."""
        if input_type == "text":
            return Path(inputfile).read_text(encoding="utf-8", errors="ignore")

        if input_type == "pdf":
            texts = []
            with pdfplumber.open(inputfile) as pdf:
                for page in pdf.pages:
                    t = page.extract_text() or ""
                    texts.append(t)
            return "\n".join(texts)

        if input_type == "word":
            doc = Document(inputfile)
            return "\n".join(p.text for p in doc.paragraphs)

        if input_type == "pptx":
            prs = Presentation(inputfile)
            slides = []
            for slide in prs.slides:
                parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        parts.append(shape.text)
                slides.append("\n".join(parts))
            return "\n\n".join(slides)

        return ""

    def extract_text_from_image(self, inputfile: str) -> str:
        """OCR via Tesseract (deu+eng)."""
        img = Image.open(inputfile)
        txt = pytesseract.image_to_string(img, lang="deu+eng")
        return txt

    def _audio_placeholder(self, inputfile: str) -> str:
        name = Path(inputfile).name
        return f"[AUDIO transcription unavailable: {name}]"

    def extract_text_from_audio(self, inputfile: str) -> str:
        """ASR via Whisper model (forced German; adjust as needed)."""
        if not self.audio_available or self.whisper_model is None:
            logger.info(
                "Audio transcription not available - using placeholder for %s.",
                _relpath(str(inputfile)),
            )
            return self._audio_placeholder(inputfile)

        result = self.whisper_model.transcribe(inputfile, language="de")
        text = (result.get("text", "") or "").strip()
        return text if text else self._audio_placeholder(inputfile)

    def extract_text_from_video(self, inputfile: str) -> str:
        """
        Simple video processing:
        - Extract audio -> temporary WAV in the current directory
        - Whisper transcription
        """
        _require_dependency(VideoFileClip, "moviepy", _MOVIEPY_IMPORT_ERROR)
        base = Path(inputfile).stem
        tmp_wav = f"{base}.veco_tmp.wav"  # relative path, avoids user home directory
        clip = VideoFileClip(inputfile)
        text = ""
        try:
            clip.audio.write_audiofile(tmp_wav, verbose=False, logger=None)
            text = self.extract_text_from_audio(tmp_wav)
        finally:
            try:
                clip.close()
            except Exception:
                pass
            try:
                if clip.audio is not None:
                    clip.audio.close()
            except Exception:
                pass
            try:
                os.remove(tmp_wav)
            except Exception:
                pass
        return text

    # ----------------------- OPTIONAL: Diarization ---------------------
    def _run_diarization(self, inputfile: str, diarization_kwargs: Optional[dict] = None) -> Optional[str]:
        """
        Runs your external diarization pipeline when available:
        - erwartet ein Modul 'veco_diarization.py' mit 'run_file' und 'build_config'
        - returns speaker-tagged text or None when the module fails/is missing
        """
        dia = _try_import_diarization()
        if dia is None:
            logger.info("Diarization module (veco_diarization.py) not found - skipping.")
            return None

        with tempfile.TemporaryDirectory(prefix="veco_dia_") as tmpdir:
            try:
                kwargs = diarization_kwargs or {}
                kwargs.setdefault("audio_dir", tmpdir)
                ok, out_txt = dia.run_file(inputfile, **kwargs)
                if ok and out_txt and os.path.exists(out_txt):
                    return Path(out_txt).read_text(encoding="utf-8", errors="ignore")
            except Exception as e:
                logger.warning(f"Diarization fehlgeschlagen: {e}")
        return None

    # ----------------------- OPTIONAL: Vision/CNN ----------------------
    def _init_vision_classifier(self):
        if self._vision_cls is not None:
            return
        if not _VISION_OK:
            logger.info("torchvision not available - skipping image classification.")
            return
        try:
            weights = torchvision.models.ResNet50_Weights.DEFAULT
            self._vision_cls = torchvision.models.resnet50(weights=weights)
            self._vision_cls.eval()
            self._vision_tf = weights.transforms()
        except Exception:
            # Fallback to simple normalisation/resize
            self._vision_cls = torchvision.models.resnet50(pretrained=True)
            self._vision_cls.eval()
            self._vision_tf = transforms.Compose([
                transforms.Resize(256),
                transforms.CenterCrop(224),
                transforms.ToTensor(),
                transforms.Normalize(mean=[0.485, 0.456, 0.406],
                                     std =[0.229, 0.224, 0.225]),
            ])

    def _image_classify(self, inputfile: str, topk: int = 5) -> Optional[str]:
        self._init_vision_classifier()
        if self._vision_cls is None or self._vision_tf is None:
            return None
        try:
            img = Image.open(inputfile).convert("RGB")
            x = self._vision_tf(img).unsqueeze(0)
            with torch.no_grad():
                logits = self._vision_cls(x)
                probs = torch.softmax(logits, dim=1)[0]
                k = min(topk, probs.shape[0])
                topv, topi = torch.topk(probs, k=k)
                try:
                    labels = torchvision.models.ResNet50_Weights.DEFAULT.meta["categories"]
                except Exception:
                    labels = [f"class_{i}" for i in range(probs.shape[0])]
                pairs = [f"{labels[int(i)]} ({float(v):.2%})" for v, i in zip(topv, topi)]
                return "Bildklassifikation (Top): " + ", ".join(pairs)
        except Exception as e:
            logger.warning(f"Bildklassifikation fehlgeschlagen: {e}")
            return None

    def _image_caption_external(self, inputfile: str, **kwargs) -> Optional[str]:
        """
        Nutzt ein externes Projekt 'veco_pic_describe' (falls vorhanden).
        Erwartete API (mindestens eine der folgenden):
          - describe(image_path: str, **kwargs) -> str
          - run(image_path: str, **kwargs) -> str
        """
        mod = _try_import_pic_describe()
        if mod is None:
            logger.info("veco_pic_describe not found - skipping image captioning.")
            return None
        try:
            if hasattr(mod, "describe") and callable(mod.describe):
                text = mod.describe(inputfile, **kwargs)
            elif hasattr(mod, "run") and callable(mod.run):
                text = mod.run(inputfile, **kwargs)
            else:
                logger.warning("veco_pic_describe hat weder describe() noch run().")
                return None
            text = (text or "").strip()
            return f"Bildbeschreibung: {text}" if text else None
        except Exception as e:
            logger.warning(f"Externe Bildbeschreibung fehlgeschlagen: {e}")
            return None

    # ----------------------- Embedding & Index -------------------------
    def embed_texts(self, texts: List[str]) -> np.ndarray:
        """Batch embedding for a list of strings (returns a float32 ndarray)."""
        if not texts:
            dim = self.embedder.get_sentence_embedding_dimension()
            return np.zeros((0, dim), dtype=np.float32)
        vecs = self.embedder.encode(texts, convert_to_numpy=True, normalize_embeddings=False)
        return np.asarray(vecs, dtype=np.float32)

    def _reserve_vector_ids(self, count: int) -> np.ndarray:
        if count <= 0:
            return np.zeros((0,), dtype=np.int64)
        start = self._next_vector_id
        self._next_vector_id += count
        return np.arange(start, start + count, dtype=np.int64)

    def _allocate_doc_id(self) -> int:
        doc_id = self._next_doc_id
        self._next_doc_id += 1
        return doc_id

    def _reset_in_memory_state(self):
        """Clear cached records, index and counters."""
        self.outputdb.clear()
        self.id_lookup.clear()
        self._next_vector_id = 0
        self._next_doc_id = 0
        if getattr(self, "_embedding_dim", None) is not None:
            self.faiss_index = IndexIDMap(IndexFlatL2(self._embedding_dim))

    def _track_existing_record(self, rec: Dict[str, Any]):
        """Register a record that was loaded from JSON or external storage."""
        rid = rec.get("id")
        if rid is not None:
            try:
                rid_int = int(rid)
            except (TypeError, ValueError):
                rid_int = None
            else:
                vec = rec.get("vector")
                if isinstance(vec, (list, tuple)):
                    try:
                        arr = np.asarray(vec, dtype=np.float32)
                    except Exception as exc:  # pragma: no cover - defensive guard
                        logger.warning(f"Skipped vector for record {rid_int}: {exc}")
                    else:
                        if arr.size > 0:
                            arr = arr.reshape(1, -1)
                            dim = getattr(self, "_embedding_dim", None)
                            if dim is not None and arr.shape[1] != dim:
                                if arr.size == dim:
                                    arr = arr.reshape(1, dim)
                                else:
                                    logger.warning(
                                        "Skipped vector for record %s: expected dim %s, got %s",
                                        rid_int,
                                        dim,
                                        arr.shape[1],
                                    )
                                    arr = None
                            if arr is not None:
                                ids = np.array([rid_int], dtype=np.int64)
                                self.faiss_index.add_with_ids(arr, ids)
                                self.id_lookup[rid_int] = rec
                if rid_int is not None and rid_int >= self._next_vector_id:
                    self._next_vector_id = rid_int + 1

        doc_id = rec.get("doc_id")
        if doc_id is not None:
            try:
                doc_int = int(doc_id)
            except (TypeError, ValueError):
                return
            if doc_int >= self._next_doc_id:
                self._next_doc_id = doc_int + 1

    def _add_records(self, vectors: np.ndarray, chunks: List[str], source: str, doc_id: int):
        """
        Insert chunks and embeddings into the index/database with consistent FAISS IDs.
        - vectors: (N, dim)
        - chunks:  list of chunk texts (N)
        - source:  relative path to the source
        - doc_id:  logical grouping (first ID for this document)
        """
        assert vectors.shape[0] == len(chunks)
        ids = self._reserve_vector_ids(len(chunks))
        if ids.size == 0:
            return

        # FAISS: add embeddings with explicit IDs
        self.faiss_index.add_with_ids(vectors, ids)

        # Build records and persist (in-memory plus optional external storage)
        src = _relpath(source)
        for local_idx, (rid, chunk, vec) in enumerate(zip(ids.tolist(), chunks, vectors)):
            rec: Dict[str, Any] = {
                "id": int(rid),
                "doc_id": int(doc_id),
                "chunk_idx": local_idx,
                "text": chunk,              # Embedding base = original chunk
                "source": src,              # relative source
                "vector": vec.tolist(),     # useful for re-index/load
            }
            self.outputdb.append(rec)
            self.id_lookup[int(rid)] = rec
            if self._ext_storage is not None and self.write_through:
                self._ext_storage.upsert(rec)

    # --------------------------- Persistence ----------------------------
    def save_database(self, json_path: Optional[str] = None):
        """
        JSON fallback: writes self.outputdb to a JSON file.
        If an external storage backend is active and write_through=False, upsert there as well.
        """
        path = json_path or self.preload_json_path
        payload = {"outputdb": self.outputdb}
        Path(path).write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")

        if self._ext_storage is not None and not self.write_through:
            for rec in self.outputdb:
                self._ext_storage.upsert(rec)

        logger.info(f"JSON saved: {_relpath(str(path))}")

    def load_database(self, json_path: Optional[str] = None):
        """
        Load the JSON fallback and populate FAISS plus id_lookup.
        If the file does not exist we start empty (no exception).
        """
        self._reset_in_memory_state()

        path = json_path or self.preload_json_path
        if not Path(path).exists():
            logger.info(f"No JSON found ({_relpath(str(path))}) - starting empty.")
            return

        data = json.loads(Path(path).read_text(encoding="utf-8"))
        vector_cnt = 0
        total_cnt = 0
        for rec in data.get("outputdb", []):
            self.outputdb.append(rec)
            total_cnt += 1
            before = self.faiss_index.ntotal
            self._track_existing_record(rec)
            if self.faiss_index.ntotal > before:
                vector_cnt += 1

        logger.info(f"JSON loaded: {vector_cnt} vectors in FAISS ({total_cnt} records).")

    def _bootstrap_from_storage(self):
        """
        Load data from external storage (SQLite/Mongo) and rebuild the index/lookup.
        JSON remains untouched.
        """
        self._reset_in_memory_state()

        vector_cnt = 0
        total_cnt = 0
        for rec in self._ext_storage.load_all():
            self.outputdb.append(rec)
            total_cnt += 1
            before = self.faiss_index.ntotal
            self._track_existing_record(rec)
            if self.faiss_index.ntotal > before:
                vector_cnt += 1

        logger.info(
            f"Storage loaded: {total_cnt} records, {vector_cnt} with embeddings."
        )

    def _switch_database(self, database: str):
        """
        Dynamischer DB-Wechsel anhand des 'database'-Strings:
          - *.json           -> JSON-Fallback
          - *.sqlite / *.db  -> SQLite (braucht storages.SqliteStorage)
          - mongodb://...    -> MongoDB (braucht storages.MongoStorage)
        Baut danach FAISS & Lookup neu auf.
        """
        db = (database or "").strip()

        # Already active? Then do nothing.
        if getattr(self, "_active_db", None) == db:
            return

        # Reset index and caches
        self._reset_in_memory_state()

        # Close external storage if necessary
        if self._ext_storage is not None:
            try:
                self._ext_storage.close()
            except Exception:
                pass
        self._ext_storage = None

        # Route based on scheme/extension
        lower = db.lower()
        stor_mod = _try_import_storages()

        if lower.endswith(".json") or lower == "":
            # JSON fallback
            self.preload_json_path = db if db else self.preload_json_path
            self.load_database(self.preload_json_path)
            self._active_db = db
            logger.info(f"Switched to JSON database: {_relpath(str(self.preload_json_path))}")
            return

        if lower.endswith(".sqlite") or lower.endswith(".db"):
            if stor_mod is None:
                raise RuntimeError("SQLite requires 'storages.py'. Please provide it.")
            self._ext_storage = stor_mod.SqliteStorage(db_path=db)
            self._bootstrap_from_storage()
            self._active_db = db
            logger.info(f"Switched to SQLite database: {_relpath(db)}")
            return

        if lower.startswith("mongodb://") or lower.startswith("mongodb+srv://"):
            if stor_mod is None:
                raise RuntimeError("MongoDB requires 'storages.py'. Please provide it.")
            self._ext_storage = stor_mod.MongoStorage(uri=db, db_name="veco_db", collection="entries")
            self._bootstrap_from_storage()
            self._active_db = db
            logger.info(f"Switched to Mongo database (uri): {db}")
            return

        raise ValueError(f"Unbekanntes database-Format: {database}")

    def close(self):
        """Cleanup (close external storage cleanly)."""
        if self._ext_storage is not None:
            try:
                self._ext_storage.close()
            except Exception as e:
                logger.warning(f"Storage close error: {e}")

    # ------------------------ LLM / Summarization ----------------------
    def build_compression_prompt(self, text: str) -> str:
        return (
            "Summarise the following text as an executive summary (5-8 bullet points).\n\n"
            f"TEXT:\n{text}\n"
        )

    def ask_llm(self, prompt: str, model: Optional[str] = None) -> str:
        if ollama is None:
            raise RuntimeError("Ollama not available.")
        m = model or self.default_model
        resp = ollama.generate(model=m, prompt=prompt)
        return (resp.get("response") or "").strip()

    # ------------------------ Ingest / Vectorization ------------------
    def vectorize(
        self,
        inputfile: str,
        use_compression: bool = False,
        model: Optional[str] = None,

        # AUTO heuristic:
        use_diarization: Optional[bool] = None,   # None = AUTO, True/False = erzwingen
        diarization_kwargs: Optional[dict] = None,

        vision_mode: Optional[str] = None,        # None = AUTO (for images), "classify" | "caption" | "both" | ""
        topk: int = 5,                            # for image classification
        pic_kwargs: Optional[dict] = None,        # optional parameters for veco_pic_describe
    ):
        """
        Voller Ingest-Pipeline-Step:
          1) Extract text (depends on file type). If use_diarization=True (audio/video only):
             -> Speaker-getaggte Transkription wird bevorzugt (wenn Modul vorhanden).
          2) For images optionally run CNN classification and/or external captioning.
          3) Chunking mit Overlap
          4) Embedding (nur Original-Chunks)
          5) Index + Persistenz
          6) Optional: Summarization (stored as metadata, never used as embedding input)
        """
        spinner = Spinner("Vectorizing input")
        spinner.start()
        try:
            input_type = self.detect_input_type(inputfile)
            logger.info(f"Detected input type: {input_type}")

            raw_text = ""

            # --- 1) Full text (includes AUTO diarization for audio/video) ---
            if use_diarization is None:
                use_diarization = (input_type in {"audio", "video"}) and (_try_import_diarization() is not None)

            if use_diarization and input_type in {"audio", "video"}:
                dia_text = self._run_diarization(inputfile, diarization_kwargs=diarization_kwargs)
                if dia_text:
                    raw_text = dia_text
                else:
                    raw_text = (self.extract_text_from_audio(inputfile)
                                if input_type == "audio" else
                                self.extract_text_from_video(inputfile))
            else:
                if input_type in {"text", "pdf", "word", "pptx"}:
                    raw_text = self.extract_text(inputfile, input_type)
                elif input_type == "image":
                    raw_text = self.extract_text_from_image(inputfile)  # OCR-Basis
                elif input_type == "audio":
                    raw_text = self.extract_text_from_audio(inputfile)
                elif input_type == "video":
                    raw_text = self.extract_text_from_video(inputfile)
                else:
                    raw_text = ""

            raw_text = (raw_text or "").strip()

            # --- 2) Vision extras (AUTO) - only relevant for images ---
            vision_extra = ""
            if input_type == "image":
                if vision_mode is None:
                    can_cls = _VISION_OK
                    can_cap = (_try_import_pic_describe() is not None)
                    if can_cls and can_cap:
                        vision_mode = "both"
                    elif can_cls:
                        vision_mode = "classify"
                    elif can_cap:
                        vision_mode = "caption"
                    else:
                        vision_mode = ""

                if vision_mode in ("classify", "both"):
                    cls = self._image_classify(inputfile, topk=topk)
                    if cls:
                        vision_extra += cls
                if vision_mode in ("caption", "both"):
                    cap = self._image_caption_external(inputfile, **(pic_kwargs or {}))
                    if cap:
                        if vision_extra:
                            vision_extra += "\n"
                        vision_extra += cap

                if vision_extra:
                    raw_text = (raw_text + "\n\n" if raw_text else "") + f"[VISION]\n{vision_extra}"

            if not raw_text:
                logger.warning("No text extracted.")
                return

            # 3) Chunking
            chunks = chunk_text(raw_text, chunk_chars=1800, overlap_chars=200)
            if not chunks:
                chunks = [raw_text]

            # 4) Embedding (original chunks only)
            vectors = self.embed_texts(chunks)

            # 5) Doc ID (used for grouping; independent from FAISS IDs)
            doc_id = self._allocate_doc_id()

            # 6) Index + DB
            self._add_records(vectors, chunks, source=str(inputfile), doc_id=doc_id)

            # 7) Optional: Summary (never used as embedding input)
            if use_compression:
                try:
                    summary = self.ask_llm(self.build_compression_prompt(raw_text), model or self.default_model)
                except Exception as e:
                    logger.warning(f"Summarization failed: {e}")
                    summary = None

                if summary:
                    meta = {
                        "id": int(10_000_000_000 + doc_id),  # large ID outside the FAISS range
                        "doc_id": int(doc_id),
                        "chunk_idx": -1,
                        "kind": "doc_summary",
                        "text": "",
                        "summary": summary,
                        "source": _relpath(str(inputfile)),
                    }
                    self.outputdb.append(meta)
                    if self._ext_storage is not None and self.write_through:
                        self._ext_storage.upsert(meta)
        finally:
            spinner.stop()

    # --------------------------- Retrieval / RAG -----------------------
    def retrieve_context(self, query: str, top_k: int = 5) -> List[Dict[str, Any]]:
        """Find the top_k most similar chunks for a query (embedding + FAISS)."""
        qv = self.embed_texts([query])
        if self.faiss_index.ntotal == 0:
            return []
        D, I = self.faiss_index.search(qv, top_k)
        hits: List[Dict[str, Any]] = []
        for rid in I[0].tolist():
            if rid == -1:
                continue
            rec = self.id_lookup.get(int(rid))
            if rec:
                hits.append(rec)
        return hits

    def query_with_context(self, question: str, top_k: int = 5, include_summary: bool = True) -> Dict[str, Any]:
        """
        Return context chunks plus (optional) document summaries - no LLM answer.
        Useful for debugging retrieval results.
        """
        ctx = self.retrieve_context(question, top_k=top_k)
        response: Dict[str, Any] = {"question": question, "contexts": ctx}
        if include_summary and ctx:
            doc_ids = list({c.get("doc_id") for c in ctx if c.get("doc_id") is not None})
            summaries = [
                rec
                for rec in self.outputdb
                if rec.get("kind") == "doc_summary" and rec.get("doc_id") in doc_ids
            ]
            if summaries:
                response["summaries"] = summaries
        return response

    def _build_rag_prompt(
        self,
        question: str,
        contexts: List[Dict[str, Any]],
        summaries: Optional[List[Dict[str, Any]]] = None,
    ) -> str:
        """Build a robust German RAG prompt for Ollama."""
        ctx_text = "\n\n".join(c.get("text", "") for c in contexts if c.get("text"))
        sum_text = ""
        if summaries:
            only = [s.get("summary", "") for s in summaries if s.get("summary")]
            if only:
                sum_text = "\n\nSUMMARY (document level):\n" + "\n".join(only)

        prompt = f"""Answer the question strictly based on the following context.
If the answer is not clearly grounded in the context, respond with "Not available in the provided context."

CONTEXT:
{ctx_text}
{sum_text}

QUESTION:
{question}

ANSWER (concise, German):
"""
        return prompt

    def query(
        self,
        database: str,
        question: str,
        llm_model: str,
        top_k: int = 5,
        include_summary: bool = True,
    ) -> Dict[str, Any]:
        """
        High-level RAG-Query (End-to-End):
          - database: Path to JSON / SQLite (.sqlite/.db) or Mongo URI (mongodb://...)
          - question: user question
          - llm_model: Ollama model name (e.g., "gemma3:12b")
          - top_k:    number of context chunks
        """
        if ollama is None:
            raise RuntimeError("Ollama not available - please install/configure it.")

        # 1) Choose/load database
        self._switch_database(database)

        # 2) Context
        contexts = self.retrieve_context(question, top_k=top_k)

        # 3) Summaries (optional)
        summaries: List[Dict[str, Any]] = []
        if include_summary and contexts:
            doc_ids = list({c.get("doc_id") for c in contexts if c.get("doc_id") is not None})
            summaries = [
                rec
                for rec in self.outputdb
                if rec.get("kind") == "doc_summary" and rec.get("doc_id") in doc_ids
            ]

        # 4) Build prompt and call the LLM
        prompt = self._build_rag_prompt(question, contexts, summaries)
        try:
            resp = ollama.generate(model=llm_model, prompt=prompt)
            answer = (resp.get("response") or "").strip()
        except Exception as e:
            raise RuntimeError(f"Ollama error: {e}")

        # 5) Return result (with source list)
        result = {
            "question": question,
            "model": llm_model,
            "answer": answer,
            "contexts": contexts,
            "sources": list({c.get("source") for c in contexts if c.get("source")}),
        }
        if include_summary and summaries:
            result["summaries"] = summaries
        return result


# ------------------------------ CLI / Demo -----------------------------------
if __name__ == "__main__":
    import argparse

    ap = argparse.ArgumentParser(description="veco-ai - Vectorize & RAG Retrieval")
    ap.add_argument("input", nargs="?", help="Datei (txt/pdf/docx/pptx/image/audio/video)")
    ap.add_argument("--compress", action="store_true", help="Store generated summary alongside the document")
    ap.add_argument("--json", default="vector_db.json", help="JSON-Fallback-Datei")
    ap.add_argument("--use-sqlite", default=None, help="Pfad zu SQLite DB (optional)")
    ap.add_argument("--use-mongo", default=None, help="Mongo URI (optional, z.B. mongodb://localhost:27017)")
    ap.add_argument("--mongo-db", default="veco_db", help="Mongo DB-Name (nur CLI-Demo)")
    ap.add_argument("--mongo-col", default="entries", help="Mongo Collection (nur CLI-Demo)")
    # Optional: control vision explicitly (empty string, classify, caption, both)
    ap.add_argument("--vision", default=None, help="Bildmodus: classify|caption|both|'' (None=AUTO)")
    ap.add_argument("--topk", type=int, default=5, help="Top-K classes for image classification")
    # Optional: control diarization explicitly (None=AUTO; true/false to force)
    ap.add_argument("--diarize", default=None, choices=["true", "false"], help="Diarization erzwingen (true/false). None=AUTO")

    args = ap.parse_args()

    # Optional external storage for ingestion (CLI demo)
    storage_kind = None
    storage_kwargs = None
    if args.use_sqlite:
        storage_kind = "sqlite"
        storage_kwargs = {"db_path": args.use_sqlite}
    elif args.use_mongo:
        storage_kind = "mongo"
        storage_kwargs = {"uri": args.use_mongo, "db_name": args.mongo_db, "collection": args.mongo_col}

    veco = Vectorize(
        preload_json_path=args.json,
        storage_kind=storage_kind,
        storage_kwargs=storage_kwargs,
        write_through=True,
    )

    # CLI -> map command line flags to None/bools
    diarize_flag: Optional[bool]
    if args.diarize is None:
        diarize_flag = None  # AUTO
    else:
        diarize_flag = (args.diarize.lower() == "true")

    if args.input:
        veco.vectorize(
            args.input,
            use_compression=args.compress,
            use_diarization=diarize_flag,   # None=AUTO
            diarization_kwargs=None,         # populate with parameters if needed
            vision_mode=args.vision,         # None=AUTO
            topk=args.topk,
        )
        veco.save_database(args.json)

        # Quick retrieval test (without LLM)
        res = veco.query_with_context("What is the document about?", top_k=5, include_summary=True)
        print(json.dumps(res, ensure_ascii=False, indent=2))
    else:
        print("No input provided. Example:")
        print("  python veco.py docs/report.pdf --compress --json vector_db.json --use-sqlite data/veco.sqlite")
        print("  python veco.py sample.wav --diarize true --json vector_db.json")
        print("  python veco.py image.jpg --vision both --json vector_db.json")

    veco.close()
